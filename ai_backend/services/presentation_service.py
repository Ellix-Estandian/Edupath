from pptx import Presentation


def create_presentation(slides, output_file):
    prs = Presentation()

    for slide_data in slides:

        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = slide_data["title"]

        body = slide.placeholders[1].text_frame

        body.clear()

        for bullet in slide_data["bullets"]:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0

    prs.save(output_file)

    return output_file